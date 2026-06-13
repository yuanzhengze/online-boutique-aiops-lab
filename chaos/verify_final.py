import os, json
from pptx import Presentation

print('=' * 60)
print('FINAL VERIFICATION')
print('=' * 60)

# 1. Check screenshots exist
screenshot_dir = r'd:\DDesk\online-boutique-aiops-lab-repo\results\screenshots'
expected = ['chaos-dashboard.png', 'grafana-podkill.png', 'grafana-cpu-memory.png',
            'review-api-response.png', 'review-degradation.png',
            'review-frontend1.png', 'review-frontend2.png',
            'review-submit1.png', 'review-submit2.png']

print('\n--- Screenshot Files ---')
for f in expected:
    path = os.path.join(screenshot_dir, f)
    exists = os.path.exists(path)
    size = os.path.getsize(path) if exists else 0
    print(f'  {f}: {"OK" if exists else "MISSING"} ({size//1024}KB)')

# 2. Check report references
report_path = r'd:\DDesk\online-boutique-aiops-lab-repo\docs\谭张锐负责部分实验报告.md'
with open(report_path, 'r', encoding='utf-8') as f:
    report = f.read()

print('\n--- Report Screenshot References ---')
refs = [
    ('chaos-dashboard.png', '截图6: 1.2节'),
    ('grafana-podkill.png', '截图4: 1.5节 PodKill'),
    ('grafana-cpu-memory.png', '截图5: 1.5节 CPU/Memory'),
    ('review-api-response.png', '截图3: 3.3节 API'),
    ('review-frontend1.png', '截图1: 3.5节 前端'),
    ('review-submit1.png', '截图7: 3.5节 提交'),
    ('review-degradation.png', '截图2: 3.6节 降级'),
]
for f, desc in refs:
    found = f in report
    print(f'  {desc}: {"OK" if found else "MISSING"}')

# 3. Check for old wrong data in report
print('\n--- Checking for stale data in report ---')
stale_patterns = ['1.71s', 'N/A(500)', '~29/s', '约 30% 请求超时或失败']
for p in stale_patterns:
    if p in report:
        print(f'  WARNING: Found stale data: "{p}"')
    else:
        print(f'  OK: No stale data "{p}"')

# 4. Check PPT
ppt_path = r'd:\DDesk\online-boutique-aiops-lab-repo\docs\谭张锐答辩PPT.pptx'
prs = Presentation(ppt_path)
print(f'\n--- PPT: {len(prs.slides)} slides ---')
for i, slide in enumerate(prs.slides):
    imgs = sum(1 for s in slide.shapes if s.shape_type == 13)
    print(f'  Slide {i+1}: {imgs} images')

# Check PPT text for stale data
print('\n--- Checking PPT for stale data ---')
all_ppt_text = ''
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                all_ppt_text += para.text + '\n'

stale_ppt = ['1.71s', 'N/A(500)', '2.21s', '~29/s', '不稳定', '>0%', '~30%',
             '丢包率与失败率基本一致', '响应时间增加 ~300ms']
for p in stale_ppt:
    if p in all_ppt_text:
        print(f'  WARNING: Stale PPT data: "{p}"')
    else:
        print(f'  OK: PPT clean of "{p}"')

# 5. Verify JSON data
print('\n--- Data Consistency Check ---')
with open(r'd:\DDesk\online-boutique-aiops-lab-repo\results\test-results\jmeter\analysis_summary.json', 'r') as f:
    jdata = json.load(f)

b50 = jdata['baseline-50u']
print(f'  JSON baseline-50u: avg={b50["avg"]}ms, P99={b50["p99"]}ms, throughput={b50["throughput"]}/s')
pk10 = jdata['pod-kill-cartservice-10u']
print(f'  JSON pod-kill-10u: errors={pk10["errors"]}, error_rate={pk10["error_rate"]}')

print()
print('=' * 60)
print('VERIFICATION COMPLETE')
print('=' * 60)